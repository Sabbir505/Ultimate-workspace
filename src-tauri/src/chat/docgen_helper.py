"""relay_docgen — editorial, frontier-quality styled document builders.

Auto-placed on the Python path for every ``generate_document`` run. Import it to
produce professionally designed DOCX, PPTX and PDF files (the kind of output
Claude / ChatGPT / Gemini ship) with a handful of high-level calls, instead of
hand-writing low-level python-docx / python-pptx / reportlab styling.

Design language (kept deliberately restrained — this is what separates a modern
deliverable from a dated "corporate template"):

* Hierarchy comes from **type scale, weight and whitespace** — never from
  decorative accent bars, stripes or underlines under titles.
* A single restrained accent colour, used sparingly. Clean white pages; rich
  full-bleed colour reserved for cover / section slides.
* An editorial **serif display** face paired with a clean **sans** body.
* Real list formatting (Word list styles / PowerPoint bullet XML), never literal
  "•" characters glued onto text.

Everything saves to the path in the ``RELAY_OUTPUT`` environment variable.

    import relay_docgen as cd
    doc = cd.Doc(title="Quarterly Review", subtitle="FY2025 · Q2",
                 theme="ink", author="Acme Analytics")
    doc.heading("Overview"); doc.paragraph("...")
    doc.bullets(["First point", "Second point"]); doc.numbered(["Step one", "Step two"])
    doc.callout("Key takeaway."); doc.table(["Metric", "Value"], [["Revenue", "$1.2M"]])
    doc.save()

    deck = cd.Deck(title="Product Launch", subtitle="2025 Roadmap",
                   theme="midnight", footer="Acme Inc")
    deck.section("Introduction", number=1)
    deck.bullets("Why now?", ["Market ready", "Tech mature"], eyebrow="Context")
    deck.two_column("Compare", "Option A", ["fast"], "Option B", ["robust"])
    deck.table_slide("Numbers", ["Q", "Rev"], [["Q1", "1.0"], ["Q2", "1.2"]])
    deck.closing("Thank you", "questions@acme.com"); deck.save()

    pdf = cd.Pdf(title="Research Brief", subtitle="On-device inference",
                 theme="plum", author="Acme Labs")
    pdf.heading("Summary"); pdf.paragraph("..."); pdf.bullets(["a", "b"])
    pdf.table(["Model", "Latency"], [["A", "12ms"]]); pdf.callout("Bottom line.")
    pdf.save()

Themes: "ink" (default), "midnight", "emerald", "plum", "amber", "crimson",
"teal". Drop to the raw objects (``doc.document`` / ``deck.prs``) for anything
not covered.
"""

import datetime
import json
import os

# ---------------------------------------------------------------------------
# Design tokens — loaded from the SHARED token file (docdesign_tokens.json,
# staged next to this module by the host) so the Python engine styles from the
# same source of truth as the JS engines and the HTML print CSS. The literals
# below are only a fallback for standalone imports without the staged file.
# Hex strings (no leading '#').
#   ink     – body text            muted   – secondary text
#   accent  – the one accent hue    tint    – faint accent wash for fills
#   hair    – hairline rules
#   cbg/cfg – cover background / foreground
#   cmut    – cover secondary text  cacc    – cover accent (title-slide eyebrow)
# ---------------------------------------------------------------------------
_FALLBACK_THEMES = {
    "ink":      {"ink": "14161C", "muted": "6B7280", "accent": "2F55E0", "tint": "EEF2FE",
                 "hair": "E5E7EB", "cbg": "0E1116", "cfg": "FFFFFF", "cmut": "9AA4B2", "cacc": "9DB6FF"},
    "midnight": {"ink": "14161C", "muted": "6B7280", "accent": "6D4AE0", "tint": "F1EEFE",
                 "hair": "E7E5EF", "cbg": "141033", "cfg": "FFFFFF", "cmut": "A6A0C8", "cacc": "BBA9FF"},
    "emerald":  {"ink": "12211A", "muted": "5E6B64", "accent": "0E9F6E", "tint": "E8F7F0",
                 "hair": "DCE7E2", "cbg": "07281E", "cfg": "FFFFFF", "cmut": "94AEA4", "cacc": "68E0B0"},
    "plum":     {"ink": "1E1524", "muted": "6B6472", "accent": "8B3FD6", "tint": "F5EDFC",
                 "hair": "E9E2EF", "cbg": "23103A", "cfg": "FFFFFF", "cmut": "B3A3C2", "cacc": "D2A8FF"},
    "amber":    {"ink": "231A12", "muted": "6E6357", "accent": "C2740A", "tint": "FBF1E2",
                 "hair": "ECE3D6", "cbg": "2A1B0C", "cfg": "FFFFFF", "cmut": "C0AB93", "cacc": "F4C27A"},
    "crimson":  {"ink": "201414", "muted": "6E5A5A", "accent": "D33C4E", "tint": "FCECEE",
                 "hair": "EDDEDE", "cbg": "1C0F12", "cfg": "FFFFFF", "cmut": "C4A5A8", "cacc": "F5A3AC"},
    "teal":     {"ink": "10201F", "muted": "5C6B6A", "accent": "0E8C9E", "tint": "E6F5F6",
                 "hair": "D9E7E7", "cbg": "062226", "cfg": "FFFFFF", "cmut": "92AEAF", "cacc": "6FD8E0"},
}


def _load_tokens():
    """Read the shared docdesign token file; fall back to embedded copies."""
    path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "docdesign_tokens.json")
    try:
        with open(path, "r", encoding="utf-8") as fh:
            return json.load(fh)
    except Exception:
        return None


_TOKENS = _load_tokens()


def _themes_from_tokens(tokens):
    themes = {}
    for tid, theme in (tokens.get("themes") or {}).items():
        c = theme.get("color") or {}
        themes[tid] = {
            "ink": c.get("ink", "14161C"),
            "muted": c.get("muted", "6B7280"),
            "accent": c.get("accent", "2F55E0"),
            "tint": c.get("tint", "EEF2FE"),
            "hair": c.get("hair", "E5E7EB"),
            "cbg": c.get("coverBg", "0E1116"),
            "cfg": c.get("coverFg", "FFFFFF"),
            "cmut": c.get("coverMuted", "9AA4B2"),
            "cacc": c.get("coverAccent", "9DB6FF"),
        }
    return themes


THEMES = _themes_from_tokens(_TOKENS) if _TOKENS else dict(_FALLBACK_THEMES)
if not THEMES:
    THEMES = dict(_FALLBACK_THEMES)

# Aliases keep older theme names working (shared file first, then legacy).
_ALIASES = {"blue": "ink", "slate": "ink", "default": "ink", "purple": "plum",
            "green": "emerald", "red": "crimson", "orange": "amber"}
if _TOKENS:
    _ALIASES.update(_TOKENS.get("aliases") or {})

# Widely-available editorial pairing: serif display + sans body. Faces also
# come from the shared tokens when they are staged.
DISPLAY = (_TOKENS or {}).get("faces", {}).get("display", {}).get("primary", "Georgia")
BODY = (_TOKENS or {}).get("faces", {}).get("body", {}).get("primary", "Calibri")
WHITE = "FFFFFF"


def _theme(name):
    key = (name or "ink").lower()
    key = _ALIASES.get(key, key)
    return THEMES.get(key, THEMES["ink"])


def _out_path(fallback):
    return os.environ.get("RELAY_OUTPUT", fallback)


def _today():
    return datetime.date.today().strftime("%B %d, %Y")


# ===========================================================================
# DOCX
# ===========================================================================
class Doc:
    def __init__(self, title="", subtitle="", theme="ink", author=""):
        from docx import Document
        from docx.shared import Pt, RGBColor, Inches

        self.document = Document()
        self.t = _theme(theme)
        self._Pt = Pt
        self._RGB = RGBColor

        section = self.document.sections[0]
        section.top_margin = Inches(1.05)
        section.bottom_margin = Inches(1.0)
        section.left_margin = Inches(1.15)
        section.right_margin = Inches(1.15)

        normal = self.document.styles["Normal"]
        normal.font.name = BODY
        normal.font.size = Pt(11)
        normal.font.color.rgb = RGBColor.from_string(self.t["ink"])
        pf = normal.paragraph_format
        pf.space_after = Pt(8)
        pf.line_spacing = 1.42

        self._footer_brand(author or title)
        if title:
            self._cover(title, subtitle, author)

    # -- low-level helpers --------------------------------------------------
    def _rgb(self, hexstr):
        return self._RGB.from_string(hexstr)

    def _el(self, tag):
        from docx.oxml import OxmlElement
        return OxmlElement(tag)

    def _q(self, name):
        from docx.oxml.ns import qn
        return qn(name)

    def _shade(self, cell, hexstr):
        shd = self._el("w:shd")
        shd.set(self._q("w:val"), "clear")
        shd.set(self._q("w:fill"), hexstr)
        cell._tc.get_or_add_tcPr().append(shd)

    def _cell_margins(self, cell, top=70, bottom=70, left=0, right=140):
        tcPr = cell._tc.get_or_add_tcPr()
        m = self._el("w:tcMar")
        for edge, val in (("top", top), ("bottom", bottom), ("start", left), ("end", right)):
            e = self._el(f"w:{edge}")
            e.set(self._q("w:w"), str(val))
            e.set(self._q("w:type"), "dxa")
            m.append(e)
        tcPr.append(m)

    def _table_borders(self, table, header_color, hair):
        """No outer/vertical lines; hairline row separators; strong header base."""
        tblPr = table._tbl.tblPr
        borders = self._el("w:tblBorders")
        spec = {
            "top": ("none", "0", "auto"),
            "left": ("none", "0", "auto"),
            "right": ("none", "0", "auto"),
            "insideV": ("none", "0", "auto"),
            "bottom": ("single", "4", hair),
            "insideH": ("single", "4", hair),
        }
        for edge, (val, sz, col) in spec.items():
            e = self._el(f"w:{edge}")
            e.set(self._q("w:val"), val)
            e.set(self._q("w:sz"), sz)
            e.set(self._q("w:space"), "0")
            e.set(self._q("w:color"), col)
            borders.append(e)
        tblPr.append(borders)

    def _cell_bottom(self, cell, hexstr, size):
        tcPr = cell._tc.get_or_add_tcPr()
        borders = self._el("w:tcBorders")
        bottom = self._el("w:bottom")
        bottom.set(self._q("w:val"), "single")
        bottom.set(self._q("w:sz"), str(size))
        bottom.set(self._q("w:space"), "0")
        bottom.set(self._q("w:color"), hexstr)
        borders.append(bottom)
        tcPr.append(borders)

    def _no_borders(self, table):
        tblPr = table._tbl.tblPr
        borders = self._el("w:tblBorders")
        for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
            e = self._el(f"w:{edge}")
            e.set(self._q("w:val"), "none")
            borders.append(e)
        tblPr.append(borders)

    def _run(self, p, text, size, color, bold=False, italic=False, font=BODY,
             spacing=None, caps=False):
        r = p.add_run(text.upper() if caps else text)
        r.bold = bold
        r.italic = italic
        r.font.name = font
        r.font.size = self._Pt(size)
        r.font.color.rgb = self._rgb(color)
        if spacing is not None:
            rPr = r._element.get_or_add_rPr()
            sp = self._el("w:spacing")
            sp.set(self._q("w:val"), str(spacing))
            rPr.append(sp)
        return r

    def _footer_brand(self, brand):
        from docx.enum.text import WD_TAB_ALIGNMENT
        from docx.shared import Inches
        Pt = self._Pt
        footer = self.document.sections[0].footer
        p = footer.paragraphs[0]
        p.text = ""
        tabs = p.paragraph_format.tab_stops
        tabs.add_tab_stop(Inches(6.2), WD_TAB_ALIGNMENT.RIGHT)
        self._run(p, (brand or ""), 8, self.t["muted"], spacing=6)
        self._run(p, "\t", 8, self.t["muted"])
        run = p.add_run()
        run.font.size = Pt(8)
        run.font.color.rgb = self._rgb(self.t["muted"])
        fld = self._el("w:fldSimple")
        fld.set(self._q("w:instr"), "PAGE")
        run._r.append(fld)

    def _cover(self, title, subtitle, author):
        Pt = self._Pt
        self.document.add_paragraph().paragraph_format.space_after = Pt(120)

        k = self.document.add_paragraph()
        k.paragraph_format.space_after = Pt(10)
        self._run(k, "Report", 10.5, self.t["accent"], bold=True, font=BODY,
                  spacing=60, caps=True)

        p = self.document.add_paragraph()
        p.paragraph_format.space_after = Pt(10)
        p.paragraph_format.line_spacing = 1.02
        self._run(p, title, 38, self.t["ink"], bold=True, font=DISPLAY)

        if subtitle:
            sp = self.document.add_paragraph()
            sp.paragraph_format.space_after = Pt(4)
            self._run(sp, subtitle, 15, self.t["muted"], font=BODY)

        meta = self.document.add_paragraph()
        meta.paragraph_format.space_before = Pt(18)
        parts = [x for x in (author, _today()) if x]
        self._run(meta, "   ·   ".join(parts), 10, self.t["muted"], spacing=10, caps=True)

        self.document.add_page_break()

    # -- public API ---------------------------------------------------------
    def heading(self, text, level=1):
        Pt = self._Pt
        p = self.document.add_paragraph()
        pf = p.paragraph_format
        pf.space_before = Pt(22 if level == 1 else 14)
        pf.space_after = Pt(6 if level == 1 else 3)
        pf.keep_with_next = True
        if level == 1:
            self._run(p, text, 18, self.t["ink"], bold=True, font=DISPLAY)
        elif level == 2:
            self._run(p, text, 13.5, self.t["ink"], bold=True, font=BODY)
        else:
            self._run(p, text, 10.5, self.t["accent"], bold=True, font=BODY,
                      spacing=30, caps=True)
        return p

    def paragraph(self, text):
        return self.document.add_paragraph(str(text))

    def bullets(self, items):
        for it in items:
            self.document.add_paragraph(str(it), style="List Bullet")

    def numbered(self, items):
        for it in items:
            self.document.add_paragraph(str(it), style="List Number")

    def callout(self, text, label="KEY POINT"):
        Pt = self._Pt
        table = self.document.add_table(rows=1, cols=1)
        table.autofit = True
        self._no_borders(table)
        cell = table.cell(0, 0)
        self._shade(cell, self.t["tint"])
        self._cell_margins(cell, top=150, bottom=150, left=190, right=190)
        cell.paragraphs[0].text = ""
        if label:
            lp = cell.paragraphs[0]
            lp.paragraph_format.space_after = Pt(4)
            self._run(lp, label, 8.5, self.t["accent"], bold=True, spacing=40, caps=True)
            bp = cell.add_paragraph()
        else:
            bp = cell.paragraphs[0]
        bp.paragraph_format.line_spacing = 1.38
        self._run(bp, text, 11.5, self.t["ink"], font=DISPLAY)
        self.document.add_paragraph().paragraph_format.space_after = Pt(4)
        return table

    def table(self, headers, rows):
        Pt = self._Pt
        table = self.document.add_table(rows=1, cols=len(headers))
        table.autofit = True
        self._table_borders(table, self.t["ink"], self.t["hair"])
        hdr = table.rows[0].cells
        for i, h in enumerate(headers):
            self._cell_margins(hdr[i])
            self._cell_bottom(hdr[i], self.t["ink"], 12)
            para = hdr[i].paragraphs[0]
            para.text = ""
            self._run(para, str(h), 9.5, self.t["ink"], bold=True, spacing=20, caps=True)
        for row in rows:
            cells = table.add_row().cells
            for ci in range(len(headers)):
                self._cell_margins(cells[ci])
                if ci < len(row):
                    para = cells[ci].paragraphs[0]
                    para.text = ""
                    self._run(para, str(row[ci]), 10.5, self.t["ink"])
        self.document.add_paragraph().paragraph_format.space_after = Pt(4)
        return table

    def page_break(self):
        self.document.add_page_break()

    def save(self, path=None):
        self.document.save(_out_path(path or "document.docx"))


# ===========================================================================
# PPTX
# ===========================================================================
class Deck:
    MX = 0.92  # left/right margin in inches

    def __init__(self, title="", subtitle="", theme="ink", footer=""):
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.t = _theme(theme)
        self.footer = footer
        self._Inches = Inches
        self._Pt = Pt
        self._Emu = Emu
        self._RGB = RGBColor
        self._PP_ALIGN = PP_ALIGN
        self._ANCHOR = MSO_ANCHOR
        self._blank = self.prs.slide_layouts[6]
        if title:
            self.title_slide(title, subtitle)

    # -- helpers ------------------------------------------------------------
    def _rgb(self, hexstr):
        return self._RGB.from_string(hexstr)

    def _slide(self, bg=WHITE):
        slide = self.prs.slides.add_slide(self._blank)
        self._rect(slide, 0, 0, self.prs.slide_width, self.prs.slide_height, bg)
        return slide

    def _rect(self, slide, x, y, w, h, fill):
        from pptx.enum.shapes import MSO_SHAPE
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = self._rgb(fill)
        shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def _text(self, slide, x, y, w, h, runs, align=None, anchor=None,
              line_spacing=None, font=None):
        box = slide.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        if anchor is not None:
            tf.vertical_anchor = anchor
        if isinstance(runs, str):
            runs = [(runs, 18, self.t["ink"], False)]
        p = tf.paragraphs[0]
        if align is not None:
            p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        for spec in runs:
            text, size, color = spec[0], spec[1], spec[2]
            bold = spec[3] if len(spec) > 3 else False
            spacing = spec[4] if len(spec) > 4 else None
            r = p.add_run()
            r.text = text
            r.font.size = self._Pt(size)
            r.font.bold = bold
            r.font.name = font or BODY
            r.font.color.rgb = self._rgb(color)
            if spacing is not None:
                r._r.get_or_add_rPr().set("spc", str(int(spacing)))
        return box

    def _bullet(self, p, char, color):
        """Attach a real bullet glyph + hanging indent to a paragraph."""
        from pptx.oxml.ns import qn
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", "292608")
        pPr.set("indent", "-292608")
        for tag in ("a:buClr", "a:buFont", "a:buChar", "a:buAutoNum", "a:buNone"):
            for e in pPr.findall(qn(tag)):
                pPr.remove(e)
        buClr = pPr.makeelement(qn("a:buClr"), {})
        srgb = pPr.makeelement(qn("a:srgbClr"), {"val": color})
        buClr.append(srgb)
        buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial", "pitchFamily": "34"})
        buChar = pPr.makeelement(qn("a:buChar"), {"char": char})
        pPr.append(buClr)
        pPr.append(buFont)
        pPr.append(buChar)

    def _eyebrow_title(self, slide, title, eyebrow=None):
        Inches = self._Inches
        top = Inches(0.62)
        if eyebrow:
            self._text(slide, Inches(self.MX), top, Inches(11), Inches(0.34),
                       [(eyebrow.upper(), 12, self.t["accent"], True, 180)])
            top = Inches(1.0)
        self._text(slide, Inches(self.MX - 0.02), top, Inches(11.8), Inches(1.0),
                   [(title, 32, self.t["ink"], True)], font=DISPLAY, line_spacing=1.02)
        return top

    def _footer(self, slide):
        Inches = self._Inches
        idx = len(self.prs.slides._sldIdLst)
        if self.footer:
            self._text(slide, Inches(self.MX), Inches(7.04), Inches(9), Inches(0.32),
                       [(self.footer, 9, self.t["muted"], False, 20)])
        self._text(slide, Inches(11.4), Inches(7.04), Inches(1.0), Inches(0.32),
                   [(str(idx).zfill(2), 9, self.t["muted"], False)],
                   align=self._PP_ALIGN.RIGHT)

    def _bullets_box(self, slide, items, top, left=None, width=None, size=18):
        Inches, Pt = self._Inches, self._Pt
        box = slide.shapes.add_textbox(left or Inches(self.MX), top,
                                       width or Inches(11.5), Inches(4.7))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = 0
        for i, it in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(14)
            p.line_spacing = 1.12
            r = p.add_run()
            r.text = str(it)
            r.font.size = Pt(size)
            r.font.name = BODY
            r.font.color.rgb = self._rgb(self.t["ink"])
            self._bullet(p, "•", self.t["accent"])
        return box

    # -- public API ---------------------------------------------------------
    def title_slide(self, title, subtitle="", eyebrow="PRESENTATION"):
        Inches = self._Inches
        slide = self._slide(self.t["cbg"])
        if eyebrow:
            self._text(slide, Inches(self.MX), Inches(2.55), Inches(11), Inches(0.4),
                       [(eyebrow.upper(), 13, self.t["cacc"], True, 240)])
        self._text(slide, Inches(self.MX - 0.02), Inches(3.05), Inches(11.6), Inches(2.0),
                   [(title, 50, self.t["cfg"], True)], font=DISPLAY, line_spacing=1.0)
        if subtitle:
            self._text(slide, Inches(self.MX), Inches(4.95), Inches(11.2), Inches(0.9),
                       [(subtitle, 20, self.t["cmut"], False)])
        self._text(slide, Inches(self.MX), Inches(6.75), Inches(11), Inches(0.4),
                   [((self.footer or _today()), 11, self.t["cmut"], False, 20)])
        return slide

    def section(self, title, number=None):
        Inches = self._Inches
        slide = self._slide(self.t["cbg"])
        if number is not None:
            self._text(slide, Inches(self.MX - 0.04), Inches(2.15), Inches(6), Inches(2.2),
                       [(str(number).zfill(2), 150, self.t["cacc"], True)], font=DISPLAY)
        self._text(slide, Inches(self.MX), Inches(4.35), Inches(11.4), Inches(1.4),
                   [(title, 40, self.t["cfg"], True)], font=DISPLAY, anchor=self._ANCHOR.TOP)
        return slide

    def bullets(self, title, items, eyebrow=None):
        slide = self._slide()
        top = self._eyebrow_title(slide, title, eyebrow)
        self._bullets_box(slide, items, self._Inches(top.inches + 1.15))
        self._footer(slide)
        return slide

    def two_column(self, title, left_head, left_items, right_head, right_items, eyebrow=None):
        Inches = self._Inches
        slide = self._slide()
        top = self._eyebrow_title(slide, title, eyebrow)
        cy = Inches(top.inches + 1.05)
        ch = Inches(6.5 - top.inches)
        for x, head, items in (
            (Inches(self.MX), left_head, left_items),
            (Inches(6.95), right_head, right_items),
        ):
            self._rect(slide, x, cy, Inches(5.45), ch, self.t["tint"])
            self._text(slide, x + Inches(0.32), cy + Inches(0.28), Inches(4.8), Inches(0.5),
                       [(head, 17, self.t["accent"], True)])
            self._bullets_box(slide, items, cy + Inches(0.95),
                              left=x + Inches(0.32), width=Inches(4.8), size=15)
        self._footer(slide)
        return slide

    def table_slide(self, title, headers, rows, eyebrow=None):
        Inches, Pt = self._Inches, self._Pt
        slide = self._slide()
        top = self._eyebrow_title(slide, title, eyebrow)
        nrows, ncols = len(rows) + 1, len(headers)
        gy = Inches(top.inches + 1.1)
        gt = slide.shapes.add_table(nrows, ncols, Inches(self.MX), gy,
                                    Inches(11.5), Inches(0.5 * nrows)).table
        gt.first_row = False
        gt.horz_banding = False
        for c, h in enumerate(headers):
            cell = gt.cell(0, c)
            cell.text = str(h).upper()
            cell.fill.solid()
            cell.fill.fore_color.rgb = self._rgb(self.t["ink"])
            r = cell.text_frame.paragraphs[0].runs[0]
            r.font.bold = True
            r.font.color.rgb = self._rgb(WHITE)
            r.font.size = Pt(13)
            r.font.name = BODY
        for ri, row in enumerate(rows, start=1):
            for c in range(ncols):
                cell = gt.cell(ri, c)
                cell.text = str(row[c]) if c < len(row) else ""
                cell.fill.solid()
                cell.fill.fore_color.rgb = self._rgb(WHITE if ri % 2 else self.t["tint"])
                if cell.text_frame.paragraphs[0].runs:
                    run = cell.text_frame.paragraphs[0].runs[0]
                    run.font.size = Pt(13)
                    run.font.name = BODY
                    run.font.color.rgb = self._rgb(self.t["ink"])
        self._footer(slide)
        return slide

    def closing(self, title, subtitle=""):
        Inches = self._Inches
        slide = self._slide(self.t["cbg"])
        self._text(slide, Inches(self.MX - 0.02), Inches(3.0), Inches(11.5), Inches(1.6),
                   [(title, 46, self.t["cfg"], True)], font=DISPLAY, anchor=self._ANCHOR.TOP)
        if subtitle:
            self._text(slide, Inches(self.MX), Inches(4.55), Inches(11.4), Inches(0.8),
                       [(subtitle, 19, self.t["cacc"], False)])
        return slide

    def save(self, path=None):
        self.prs.save(_out_path(path or "deck.pptx"))


# ===========================================================================
# PDF (reportlab) — fonts embedded for portability.
# ===========================================================================
_PDF_FONTS = None


def _register_pdf_fonts():
    """Register an embedded serif+sans family; fall back to built-ins.

    Returns (display, display_bold, body, body_bold, body_italic) font names.
    """
    global _PDF_FONTS
    if _PDF_FONTS is not None:
        return _PDF_FONTS
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    # Cross-platform font roots. Non-existent dirs are skipped, so listing all
    # platforms here is harmless. Windows fonts can live under the system dir
    # (%WINDIR%\Fonts) or per-user (%LOCALAPPDATA%\Microsoft\Windows\Fonts,
    # the default target for user-installed fonts since Windows 10).
    win_dir = os.environ.get("WINDIR", "C:\\Windows")
    local_app = os.environ.get("LOCALAPPDATA", "")
    search = [
        # Linux
        "/usr/share/fonts", "/usr/local/share/fonts",
        os.path.expanduser("~/.fonts"), os.path.expanduser("~/.local/share/fonts"),
        # Windows (system + per-user)
        os.path.join(win_dir, "Fonts"),
        os.path.join(local_app, "Microsoft", "Windows", "Fonts") if local_app else "",
        # macOS (system + per-user)
        "/Library/Fonts", "/System/Library/Fonts",
        os.path.expanduser("~/Library/Fonts"),
    ]
    search = [s for s in search if s]

    def find(names):
        low = [n.lower() for n in names]
        for root in search:
            if not os.path.isdir(root):
                continue
            for dirpath, _dirs, files in os.walk(root):
                for f in files:
                    if f.lower() in low:
                        return os.path.join(dirpath, f)
        return None

    def reg(alias, names):
        path = find(names)
        if not path:
            return None
        try:
            pdfmetrics.registerFont(TTFont(alias, path))
            return alias
        except Exception:
            return None

    disp = reg("CDDisplay", ["Georgia.ttf", "LiberationSerif-Regular.ttf",
                             "DejaVuSerif.ttf", "NotoSerif-Regular.ttf"])
    disp_b = reg("CDDisplayB", ["Georgiab.ttf", "LiberationSerif-Bold.ttf",
                                "DejaVuSerif-Bold.ttf", "NotoSerif-Bold.ttf"])
    body = reg("CDBody", ["Calibri.ttf", "LiberationSans-Regular.ttf",
                          "DejaVuSans.ttf", "Arial.ttf", "NotoSans-Regular.ttf"])
    body_b = reg("CDBodyB", ["Calibrib.ttf", "LiberationSans-Bold.ttf",
                             "DejaVuSans-Bold.ttf", "Arialbd.ttf", "NotoSans-Bold.ttf"])
    body_i = reg("CDBodyI", ["Calibrii.ttf", "LiberationSans-Italic.ttf",
                             "DejaVuSans-Oblique.ttf", "Ariali.ttf"])

    _PDF_FONTS = (
        disp or "Times-Bold",
        disp_b or disp or "Times-Bold",
        body or "Helvetica",
        body_b or "Helvetica-Bold",
        body_i or "Helvetica-Oblique",
    )
    return _PDF_FONTS


class Pdf:
    def __init__(self, title="", subtitle="", theme="ink", author=""):
        from reportlab.lib.colors import HexColor
        from reportlab.lib.pagesizes import LETTER

        self.t = _theme(theme)
        self._HexColor = HexColor
        self._pagesize = LETTER
        self.title = title
        self.subtitle = subtitle
        self.author = author
        self.story = []
        (self.f_disp, self.f_disp_b, self.f_body,
         self.f_body_b, self.f_body_i) = _register_pdf_fonts()
        self._styles()
        if title:
            self._cover(title, subtitle, author)

    def _c(self, hexstr):
        return self._HexColor("#" + hexstr)

    def _styles(self):
        from reportlab.lib.styles import ParagraphStyle
        from reportlab.lib.enums import TA_LEFT
        t = self.t
        self.st_kicker = ParagraphStyle("kicker", fontName=self.f_body_b, fontSize=10,
                                        textColor=self._c(t["accent"]), spaceAfter=8,
                                        leading=13)
        self.st_title = ParagraphStyle("title", fontName=self.f_disp_b, fontSize=32,
                                       textColor=self._c(t["ink"]), spaceAfter=8, leading=36)
        self.st_subtitle = ParagraphStyle("subtitle", fontName=self.f_body, fontSize=14,
                                          textColor=self._c(t["muted"]), spaceAfter=10, leading=19)
        self.st_meta = ParagraphStyle("meta", fontName=self.f_body, fontSize=9.5,
                                      textColor=self._c(t["muted"]), spaceAfter=18, leading=12)
        self.st_h1 = ParagraphStyle("h1", fontName=self.f_disp_b, fontSize=17,
                                    textColor=self._c(t["ink"]), spaceBefore=18,
                                    spaceAfter=6, leading=21)
        self.st_h2 = ParagraphStyle("h2", fontName=self.f_body_b, fontSize=12.5,
                                    textColor=self._c(t["ink"]), spaceBefore=12,
                                    spaceAfter=3, leading=16)
        self.st_body = ParagraphStyle("body", fontName=self.f_body, fontSize=10.5,
                                      textColor=self._c(t["ink"]), spaceAfter=8,
                                      leading=15.5, alignment=TA_LEFT)
        self.st_bullet = ParagraphStyle("bullet", parent=self.st_body,
                                        leftIndent=18, bulletIndent=3, spaceAfter=5)
        self.st_callout = ParagraphStyle("callout", fontName=self.f_disp, fontSize=11.5,
                                         textColor=self._c(t["ink"]), leading=16)
        self.st_call_label = ParagraphStyle("call_label", fontName=self.f_body_b, fontSize=8,
                                            textColor=self._c(t["accent"]), leading=12,
                                            spaceAfter=4)

    def _cover(self, title, subtitle, author):
        from reportlab.platypus import Paragraph, Spacer
        self.story.append(Spacer(1, 120))
        self.story.append(Paragraph("REPORT", self.st_kicker))
        self.story.append(Paragraph(title, self.st_title))
        if subtitle:
            self.story.append(Paragraph(subtitle, self.st_subtitle))
        self.story.append(Spacer(1, 10))
        parts = [x for x in (author, _today()) if x]
        self.story.append(Paragraph("&nbsp;&nbsp;·&nbsp;&nbsp;".join(parts), self.st_meta))
        self.story.append(Spacer(1, 8))

    # -- public API ---------------------------------------------------------
    def heading(self, text, level=1):
        from reportlab.platypus import Paragraph
        self.story.append(Paragraph(text, self.st_h1 if level == 1 else self.st_h2))

    def paragraph(self, text):
        from reportlab.platypus import Paragraph
        self.story.append(Paragraph(str(text), self.st_body))

    def bullets(self, items):
        from reportlab.platypus import Paragraph
        for it in items:
            self.story.append(Paragraph(str(it), self.st_bullet, bulletText="•"))

    def numbered(self, items):
        from reportlab.platypus import Paragraph
        for i, it in enumerate(items, 1):
            self.story.append(Paragraph(str(it), self.st_bullet, bulletText=f"{i}."))

    def callout(self, text, label="KEY POINT"):
        from reportlab.platypus import Paragraph, Table, TableStyle, Spacer
        from reportlab.lib.units import inch
        inner = []
        if label:
            inner.append(Paragraph(label.upper(), self.st_call_label))
        inner.append(Paragraph(text, self.st_callout))
        tbl = Table([[inner]], colWidths=[6.3 * inch])
        tbl.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, -1), self._c(self.t["tint"])),
            ("LEFTPADDING", (0, 0), (-1, -1), 18),
            ("RIGHTPADDING", (0, 0), (-1, -1), 18),
            ("TOPPADDING", (0, 0), (-1, -1), 14),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 14),
        ]))
        self.story.append(tbl)
        self.story.append(Spacer(1, 10))

    def table(self, headers, rows):
        from reportlab.platypus import Table, TableStyle, Spacer
        data = [[str(h).upper() for h in headers]] + [list(map(str, r)) for r in rows]
        tbl = Table(data, repeatRows=1, hAlign="LEFT")
        style = [
            ("TEXTCOLOR", (0, 0), (-1, 0), self._c(self.t["ink"])),
            ("FONTNAME", (0, 0), (-1, 0), self.f_body_b),
            ("FONTNAME", (0, 1), (-1, -1), self.f_body),
            ("FONTSIZE", (0, 0), (-1, 0), 8.5),
            ("FONTSIZE", (0, 1), (-1, -1), 10),
            ("TEXTCOLOR", (0, 1), (-1, -1), self._c(self.t["ink"])),
            ("LINEBELOW", (0, 0), (-1, 0), 1.1, self._c(self.t["ink"])),
            ("LINEBELOW", (0, 1), (-1, -2), 0.5, self._c(self.t["hair"])),
            ("TOPPADDING", (0, 0), (-1, -1), 8),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
            ("LEFTPADDING", (0, 0), (0, -1), 0),
            ("RIGHTPADDING", (-1, 0), (-1, -1), 0),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ]
        tbl.setStyle(TableStyle(style))
        self.story.append(tbl)
        self.story.append(Spacer(1, 10))

    def save(self, path=None):
        from reportlab.platypus import SimpleDocTemplate
        from reportlab.lib.pagesizes import LETTER
        out = _out_path(path or "document.pdf")
        doc = SimpleDocTemplate(out, pagesize=LETTER,
                                topMargin=60, bottomMargin=58,
                                leftMargin=64, rightMargin=64,
                                title=self.title, author=self.author)
        theme = self.t
        HexColor = self._HexColor
        body_font = self.f_body

        def footer(canvas, d):
            canvas.saveState()
            canvas.setFont(body_font, 8)
            canvas.setFillColor(HexColor("#" + theme["muted"]))
            if self.author or self.title:
                canvas.drawString(64, 34, self.author or self.title)
            canvas.drawRightString(LETTER[0] - 64, 34, str(d.page).zfill(2))
            canvas.restoreState()

        doc.build(self.story, onFirstPage=footer, onLaterPages=footer)
