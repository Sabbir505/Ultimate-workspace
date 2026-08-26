import os
import random
from datetime import datetime

os.makedirs("docs", exist_ok=True)
random.seed(42)

COMPANY = "NovaWorks Analytics"
TITLE = "Q3 2026 Business & Product Report"
MONTHS = ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep"]
PRODUCTS = ["Atlas Core", "Nimbus API", "Pulse Insights", "Orbit Mobile", "Zenith Cloud"]

revenue = [round(random.uniform(80, 160), 1) for _ in MONTHS]
users = [random.randint(5000, 50000) for _ in MONTHS]
churn = [round(random.uniform(1.5, 6.0), 1) for _ in MONTHS]

def txt(profit=120.0):
    return random.randint(400, 2000)

# ------------- Markdown doc -------------
md = f"""# {TITLE}

**Company:** {COMPANY}  
**Date:** {datetime.now().strftime('%Y-%m-%d')}  
**Authored by:** Data Intelligence Unit

## 1. Executive Summary

{COMPANY} delivered a strong quarter. Revenue grew steadily,
active users crossed new milestones, and product expansion continues across
four major product lines. Performance highlights are summarized in the tables below.

## 2. Key Metrics

| Metric | Latest Month ({MONTHS[-1]}) | Quarter Trend |
|--------|-----------------------------|---------------|
| Monthly Revenue (USD k) | {revenue[-1]} | {revenue[0]} -> {revenue[-1]} |
| Active Users | {users[-1]:,} | {users[0]:,} -> {users[-1]:,} |
| Churn Rate (%) | {churn[-1]} | {churn[0]}% -> {churn[-1]}% |

## 3. Revenue Trend

| Month | Revenue (k) | Users | Churn (%) |
|-------|-------------|-------|-----------|
"""
for m, r, u, c in zip(MONTHS, revenue, users, churn):
    md += f"| {m} | {r} | {u:,} | {c} |\n"

md += f"""
## 4. Product Breakdown

| Product | Team Size | Sprint Points | Releases |
|---------|-----------|---------------|----------|
"""
for p in PRODUCTS:
    md += f"| {p} | {random.randint(6, 40)} | {random.randint(120, 400)} | {random.randint(2, 8)} |\n"

md += f"""
## 5. Risks & Opportunities

- Onboarding friction remains the top churn driver.
- Expansion into new regions is budgeted for {datetime.now().year + 1}.
- Infrastructure cost per request has dropped 18% quarter-over-quarter.

## 6. Named Outcomes

1. **Revenue growth** of +{(revenue[-1]-revenue[0])/revenue[0]*100:.0f}% over the quarter.
2. **User base** increased {users[-1]-users[0]:,} net accounts.
3. **Churn reduced** by {churn[0]-churn[-1]:.1f} percentage points.
4. Three flagship products shipped GA releases.

---
*Generated automatically with randomized sample data for demonstration purposes.*
"""
with open("docs/report.md", "w", encoding="utf-8") as f:
    f.write(md)

print("markdown done")

# ------------- PowerPoint -------------
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLUE = RGBColor(0x1F, 0x4E, 0x79)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

def title_slide(title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bar = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.5), Inches(2.6), Inches(12), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = BLUE
    st = s.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12), Inches(1))
    sp = st.text_frame.paragraphs[0]; sp.text = subtitle
    sp.font.size = Pt(24); sp.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    return s

def content_slide(title, body):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bar = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.9))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(10), Inches(0.7))
    p = tb.text_frame.paragraphs[0]; p.text = title
    p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = WHITE
    box = s.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12), Inches(5.5))
    tf = box.text_frame; tf.word_wrap = True
    for i, line in enumerate(body.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(20); p.space_after = Pt(12)
    return s

title_slide(TITLE, f"{COMPANY}  |  Quarterly Performance Snapshot")

def chart_slide(title, chart_type, categories, series_dict, y_title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bar = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.9))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.7))
    p = tb.text_frame.paragraphs[0]; p.text = title
    p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = WHITE
    cdata = CategoryChartData()
    cdata.categories = categories
    for name, vals in series_dict.items():
        cdata.add_series(name, vals)
    gframe = s.shapes.add_chart(chart_type, Inches(1), Inches(1.2), Inches(11), Inches(5.7), cdata)
    return s

chart_slide("Revenue Trend (USD thousands)", XL_CHART_TYPE.COLUMN_CLUSTERED, MONTHS,
            {"Revenue": revenue}, "Revenue")
chart_slide("Active Users by Month", XL_CHART_TYPE.LINE, MONTHS,
            {"Active Users": [u / 1000 for u in users]}, "Users (thousands)")
chart_slide("Churn Rate (%)", XL_CHART_TYPE.BAR_CLUSTERED, MONTHS,
            {"Churn": churn}, "Churn %")

content_slide("Product Portfolio", "\n".join(f"• {p} — {random.randint(120,400)} sprint points, {random.randint(2,8)} releases this quarter" for p in PRODUCTS))
content_slide("Outlook & Priorities",
    "• Expand to 3 new markets by {0}\n• Reduce onboarding friction to cut churn\n• Ship AI-assisted insights for Pulse\n• Achieve 99.99% platform uptime".format(datetime.now().year + 1))
content_slide("Contact", "Data Intelligence Unit\n{company} · analytics@novaworks.example\nAll figures are illustrative sample data.".replace("{company}", COMPANY))

prs.save("docs/report.pptx")
print("ppt done")

# ------------- PDF -------------
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image)
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER

styles = getSampleStyleSheet()
title_style = ParagraphStyle("t", parent=styles["Title"], textColor=colors.HexColor("#1F4E79"), fontSize=26)
h2 = ParagraphStyle("h2", parent=styles["Heading2"], textColor=colors.HexColor("#1F4E79"))
h3 = ParagraphStyle("h3", parent=styles["Heading3"], textColor=colors.HexColor("#1F4E79"))

def pdf_chart():
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    fig, ax = plt.subplots(figsize=(7, 3.5))
    ax.plot(MONTHS, revenue, marker="o", color="#1F4E79", linewidth=2)
    ax.set_title("Monthly Revenue Trend (USD k)")
    ax.set_ylabel("Revenue")
    ax.grid(alpha=0.3)
    fig.tight_layout()
    path = os.path.join("docs", "chart.png")
    fig.savefig(path, dpi=100)
    plt.close(fig)
    return path

from reportlab.platypus import Image
doc = SimpleDocTemplate("docs/report.pdf", pagesize=A4)
story = []
story.append(Paragraph(TITLE, title_style))
story.append(Spacer(1, 0.2 * inch))
story.append(Paragraph(f"<b>{COMPANY}</b><br/>Generated {datetime.now().strftime('%Y-%m-%d')}", h3))
story.append(Spacer(1, 0.3 * inch))

story.append(Paragraph("1. Executive Summary", h2))
story.append(Paragraph(f"{COMPANY} posted another solid quarter with revenue of {revenue[-1]:.1f}k in {MONTHS[-1]} "
                       f"and a total of {users[-1]:,} active users. The following pages summarize key metrics, "
                       "trends, and the product breakdown.", styles["BodyText"]))
story.append(Spacer(1, 0.2 * inch))

story.append(Paragraph("2. Key Metrics", h2))
header = ["Metric", f"Latest ({MONTHS[-1]})"]
rows = [["Monthly Revenue (USD k)", str(revenue[-1])],
        ["Active Users", f"{users[-1]:,}"],
        ["Churn Rate (%)", str(churn[-1])]]
table = Table([header] + rows)
table.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1F4E79")),
                           ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                           ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                           ("FONTSIZE", (0, 0), (-1, -1), 10),
                           ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#eef2f7")])]))
story.append(table)
story.append(Spacer(1, 0.3 * inch))

pdf_chart()
story.append(Image("docs/chart.png", width=6.5 * inch, height=3.2 * inch))
story.append(Spacer(1, 0.3 * inch))

story.append(Paragraph("3. Monthly Breakdown", h2))
hdr = ["Month"] + list(MONTHS)
rev_row = ["Revenue (k)"] + [str(r) for r in revenue]
usr_row = ["Users"] + [f"{u:,}" for u in users]
t3 = Table([hdr, rev_row, usr_row])
t3.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1F4E79")),
                        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                        ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                        ("FONTSIZE", (0, 0), (-1, -1), 8)]))
story.append(t3)
story.append(Spacer(1, 0.3 * inch))

story.append(Paragraph("4. Product Breakdown", h2))
prod_rows = [["Product", "Team", "Sprint Points", "Releases"]]
for p in PRODUCTS:
    prod_rows.append([p, str(random.randint(6, 40)), str(random.randint(120, 400)), str(random.randint(2, 8))])
t4 = Table(prod_rows, colWidths=[2 * inch, 1 * inch, 1.4 * inch, 1.2 * inch])
t4.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1F4E79")),
                        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                        ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                        ("FONTSIZE", (0, 0), (-1, -1), 9)]))
story.append(t4)
story.append(Spacer(1, 0.3 * inch))
story.append(Paragraph("<i>All figures are randomized sample data for demonstration.</i>", styles["Italic"]))

doc.build(story)
print("pdf done")